# -*- coding:utf-8 -*-
import urllib
from bs4 import BeautifulSoup
from textrank4zh import TextRank4Keyword, TextRank4Sentence
from pptx import Presentation
from pptx.util import Inches, Pt

# 参考资料：
# https://github.com/letiantian/TextRank4ZH

# 注意，这个类主要处理中文网页的，部分纯字母数字内容会被删除
class page2ppt(object):
    def __init__(self, url, ppt_pages=5, ppt_path="./gen.ppt", content_path="./content.txt"):
        """
        :param url: 指定网页的网址
        :param ppt_pages: 目标生成ppt的页数(上限)
        :param ppt_path: 生成ppt的路径
        :param content_path: 生成的中间纯文本文件的路径
        """
        self.url = url
        self.ppt_pages = ppt_pages
        self.content_path = content_path
        self.ppt_path = ppt_path


    def get_content(self):
        """ 把指定网页处理成比较干净的文本格式，并保存到content.txt中 """

        # 读取指定url的网页
        url_content = urllib.urlopen(self.url)
        # 解析下载的网页
        soup = BeautifulSoup(url_content, "html5lib")
        fd = open(self.content_path, "w")
        # 删去脚本
        for script in soup('script'):
            script.extract()
        # 获取原始文本内容
        x = soup.get_text().encode('utf-8')
        # 逐行展开，形成待处理的句子列表
        x = x.split("\n")
        new = []
        for line in x:
            line = line.strip()  # 删除每行首尾的空白符
            line = ''.join(line.split())  # 删去中间的空白符
            if line.__len__() > 3 and not line.isalnum():
                new.append(line + "\n")
        fd.writelines(new)
        fd.close()

    def select_content(self):
        """ 根据纯文本内容生成一个字典数据结构，作为最终的ppt内容 """

        fd = open(self.content_path, "r")
        text = fd.read()
        fd.close()

        # Keyword part
        tr4w = TextRank4Keyword()
        tr4w.analyze(text=text, lower=True, window=2)
        keywords_list = tr4w.get_keywords(20, word_min_len=1)

        # Sentence part
        tr4s = TextRank4Sentence()
        tr4s.analyze(text=text, lower=True, source='all_filters')
        key_sentences_list = tr4s.get_key_sentences(num=self.ppt_pages * 10)

        # naive selection algorithm
        # 选出top self.ppt_pages的关键词作为每一页ppt的标题
        # 在key_sentences_list句子集合中选出包含该关键词的句子，至多4句
        content = {}
        for i in range(self.ppt_pages):
            keyword = keywords_list[i].word
            targets = []
            for k in key_sentences_list:
                if keyword in k.sentence:
                    targets.append(k.sentence)
            if targets.__len__() > 0 :
                if targets.__len__() > 4: # 每页至多选出4个句子
                    targets = targets[:4]
                content[keyword] = targets
        return content

    def convert2ppt(self):
        self.get_content()
        # print "get_content done"
        content = self.select_content()
        # print "select_content done"
        # 新建一个ppt实例
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        for key in content:
            # print "key =", key
            # 添加一张新的幻灯片
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            left = top = width = height = Inches(1)
            title = slide.shapes.title
            title.text = key
            # 添加一个新的文本框
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            # 小字-句子
            for sens in content[key]:
                p = tf.add_paragraph()
                p.text = sens
                p.font.size = Pt(25)
        # 保存ppt到磁盘文件
        prs.save(self.ppt_path)

